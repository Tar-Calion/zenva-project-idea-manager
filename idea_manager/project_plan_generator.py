from idea_manager.gemini_pro_client import GeminiProClient

import openpyxl
from pptx import Presentation


class ProjectPlanGenerator:

    def __init__(self, project_charter, file_prefix):
        self.project_charter = project_charter
        self.file_prefix = file_prefix

    def generate(self):

        # Generate project plan
        project_plan = self.__generate_project_plan(self.project_charter)

        # Save project plan
        self.__save_files(project_plan)

    def __generate_project_plan(self, project_charter):
        prompt = (f"Generate a project plan as a table given the following project charter:\n"
                  f"BEGIN\n"
                  f"{project_charter}\n"
                  f"END\n"
                  f"Generate a table with following columns. Fill the table with tasks based on the project charter above.\n"
                  f"-Task name\n"
                  f"-Duration\n"
                  f"-Dependencies\n"
                  f"-Status\n"
                  f"-Resources")

        print("Prompt:", prompt)

        client = GeminiProClient()
        project_plan = client.generate_output(prompt)

        print("Project plan:", project_plan)

        return project_plan

    def __save_files(self, project_plan):
        file_name = f"{self.file_prefix}_project_plan"

        parsed_project_plan = self.__parse_markdown_table(project_plan)

        self.__save_excel(parsed_project_plan, file_name)
        self.__save_powerpoint(parsed_project_plan, file_name)

    def __parse_markdown_table(self, markdown_table: str):
        lines = markdown_table.split("\n")
        table = []
        for line in lines[0:1] + lines[2:]:
            if not line.strip():
                continue
            cells = line.split("|")
            cells = [cell.replace("*", "").strip() for cell in cells]

            # remove cells on start and end if they are empty
            if cells[0] == "":
                cells = cells[1:]
            if cells[-1] == "":
                cells = cells[:-1]
            table.append(cells)
        return table

    def __save_excel(self, parsed_project_plan, file_name):

        workbook = openpyxl.Workbook()
        worksheet = workbook.active
        worksheet.title = "Project Plan"

        for row in parsed_project_plan:
            worksheet.append(row)

        workbook.save(f"output/{file_name}.xlsx")

    def __save_powerpoint(self, parsed_project_plan: list, file_name):

        presentation = Presentation()
        layout = presentation.slide_layouts[1]

        for row in parsed_project_plan[1:]:
            slide = presentation.slides.add_slide(layout)
            slide.placeholders[0].text = row[0]
            text = (f"Duration: {row[1]}\n"
                    f"Dependencies: {row[2]}\n"
                    f"Status: {row[3]}\n"
                    f"Resources: {row[4]}")

            slide.placeholders[1].text = text

        presentation.save(f"output/{file_name}.pptx")
